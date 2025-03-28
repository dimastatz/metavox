import os
from pptx import Presentation
from PIL import Image
from reportlab.pdfgen import canvas

def pptx_to_pdf(input_pptx_path, output_pdf_path):
    """
    Converts a PowerPoint presentation (.pptx) to a PDF file.

    Args:
        input_pptx_path (str): Path to the input .pptx file.
        output_pdf_path (str): Path to save the output .pdf file.
    """
    if not os.path.exists(input_pptx_path):
        raise FileNotFoundError(f"Input file not found: {input_pptx_path}")

    # Load the presentation
    presentation = Presentation(input_pptx_path)
    slide_images = []

    # Convert each slide to an image
    for i, slide in enumerate(presentation.slides):
        slide_image_path = f"slide_{i + 1}.png"
        slide.shapes._spTree.write(slide_image_path)  # Save slide as an image
        slide_images.append(slide_image_path)

    # Create a PDF from the images
    c = canvas.Canvas(output_pdf_path)
    for image_path in slide_images:
        with Image.open(image_path) as img:
            c.setPageSize(img.size)
            c.drawImage(image_path, 0, 0, width=img.width, height=img.height)
            c.showPage()
        os.remove(image_path)  # Clean up the temporary image file
    c.save()

# Example usage
if __name__ == "__main__":
    input_pptx = "example.pptx"  # Replace with your .pptx file path
    output_pdf = "example.pdf"  # Replace with your desired .pdf file path
    pptx_to_pdf(input_pptx, output_pdf)
