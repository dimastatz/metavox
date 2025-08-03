"""Document class for handling the document file"""

import os

import subprocess
from io import BytesIO

import soundfile as sf

from kokoro import KPipeline
from pptx import Presentation


def read_presentation(file_path) -> Presentation:
    """Read the presentation from the document file"""
    with open(file_path, "rb") as file:
        source_stream = BytesIO(file.read())
        return Presentation(source_stream)


def get_speaker_notes(presentation: Presentation, slide_index: int) -> str:
    """Get the speaker notes from the slide"""
    slide = presentation.slides[slide_index]
    notes_slide = slide.notes_slide
    return notes_slide.notes_text_frame.text


def get_libreoffice_version() -> str:
    """Get the version of LibreOffice installed on the system"""
    return subprocess.run(
        ["libreoffice", "--version"],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        check=True,
    )


def convert_pptx_to_pdf(input_path, output_dir) -> str:
    """Convert a PPTX file to PDF using LibreOffice"""
    _, file_name = os.path.split(input_path)
    file_name = os.path.splitext(file_name)[0] + ".pdf"
    output_path = os.path.join(output_dir, file_name)

    command = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        output_dir,
        input_path,
    ]

    subprocess.run(command, check=True)
    return output_path


def speaker_notes_to_audio(notes: str, lang_code="a") -> BytesIO:
    """Convert speaker notes to audio using TTS and return as memory stream"""
    pipeline = KPipeline(lang_code=lang_code)
    generator = pipeline(notes, voice="af_heart")

    for _, _, audio in generator:
        audio_stream = BytesIO()
        sf.write(audio_stream, audio, 24000, format="WAV")
        audio_stream.seek(0)
        return audio_stream
